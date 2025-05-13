from PyPDF2 import PdfReader
from fastapi import FastAPI, Request, status, File, UploadFile, HTTPException, Form
from pydantic import BaseModel, Field, validator
from fastapi.middleware.cors import CORSMiddleware
from fastapi.exceptions import RequestValidationError
from fastapi.responses import JSONResponse, StreamingResponse
from io import BytesIO
import sqlite3
from contextlib import asynccontextmanager
import logging
from fastapi.security import HTTPBasicCredentials
from PyPDF2.errors import PdfReadError
import pytesseract
from PIL import ImageDraw
from pdf2image import convert_from_bytes
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from ner import redact_entities1,redact_entities2,redact_entities3  # Importing ner module
from regex import redact_all # Importing regex module
from docx import Document
from pptx import Presentation
import pandas as pd


import bcrypt  # For securely comparing passwords (recommended)

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

DATABASE = "form_data.db"

# Initialize the database
def init_db():
    conn = sqlite3.connect(DATABASE)
    cursor = conn.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            first_name TEXT NOT NULL,
            last_name TEXT NOT NULL,
            username TEXT UNIQUE NOT NULL,
            password TEXT NOT NULL,
            confirm_password TEXT NOT NULL,
            security_question TEXT NOT NULL,
            security_answer TEXT NOT NULL
        )
    ''')
    conn.commit()
    conn.close()

# Lifespan context manager for startup and cleanup
@asynccontextmanager
async def lifespan(app: FastAPI):
    init_db()
    yield  # The app runs between yield and the end of this block

app = FastAPI(lifespan=lifespan)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Change to specific origins for production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Pydantic models
class User(BaseModel):
    first_name: str = Field(..., min_length=1, max_length=50)
    last_name: str = Field(..., min_length=1, max_length=50)
    username: str = Field(..., min_length=3, max_length=20)
    password: str = Field(..., min_length=6)
    confirm_password: str = Field(..., min_length=6)
    security_question: str = Field(..., min_length=1)
    security_answer: str = Field(..., min_length=1)

    @validator("confirm_password")
    def passwords_match(cls, confirm_password, values):
        if "password" in values and confirm_password != values["password"]:
            raise ValueError("Passwords do not match")
        return confirm_password

class RedactRequest(BaseModel):
    text: str
    redaction_level: int

# Root endpoint
@app.get("/")
def read_root():
    return {"message": "Welcome to the Redact API"}

# User registration
@app.post("/submit-form", status_code=status.HTTP_201_CREATED)
async def submit_form(user: User):
    try:
        conn = sqlite3.connect(DATABASE)
        cursor = conn.cursor()
        cursor.execute('''
            INSERT INTO users (first_name, last_name, username, password, confirm_password, security_question, security_answer)
            VALUES (?, ?, ?, ?, ?, ?, ?)
        ''', (
            user.first_name,
            user.last_name,
            user.username,
            user.password,
            user.confirm_password,
            user.security_question,
            user.security_answer,
        ))
        conn.commit()
        conn.close()
        return {"message": "User data stored successfully"}
    except sqlite3.IntegrityError:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Username already exists"
        )

# Get all users
@app.get("/users", status_code=status.HTTP_200_OK)
async def get_users():
    conn = sqlite3.connect(DATABASE)
    cursor = conn.cursor()
    cursor.execute("SELECT id, first_name, last_name, username, security_question FROM users")
    users = cursor.fetchall()
    conn.close()

    if not users:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="No users found")

    return [
        {
            "id": user[0],
            "first_name": user[1],
            "last_name": user[2],
            "username": user[3],
            "security_question": user[4],
        }
        for user in users
    ]



@app.post("/login", status_code=status.HTTP_200_OK)
async def login(credentials: HTTPBasicCredentials):
    """
    Endpoint to validate username and password for login.
    """
    try:
        conn = sqlite3.connect(DATABASE)
        cursor = conn.cursor()
        cursor.execute("SELECT password FROM users WHERE username = ?", (credentials.username,))
        result = cursor.fetchone()
        conn.close()

        if result:
            stored_password = result[0]
            # Compare plain-text passwords (not recommended for production)
            if credentials.password == stored_password:
                return {"message": "Login successful"}
            else:
                raise HTTPException(
                    status_code=status.HTTP_401_UNAUTHORIZED,
                    detail="Invalid password"
                )
        else:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Username not found"
            )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"An error occurred: {str(e)}"
        )


@app.post("/redact")
async def redact_text(request: dict):
    """
    API endpoint to redact sensitive information using NER and regex-based methods.
    Supports different levels of redaction:
    0 - No redaction
    1 - regex-based redaction
    2 - NER-based redaction (Method 1)
    3 - NER-based redaction (Method 2)
    4 - NER-based redaction (Method 3)
    Anything else - Returns proper error.
    """
    text = request.get("text", "")
    redaction_level = request.get("redaction_level", 0)

    try:
        if redaction_level == 0:
            return {"redacted_text": text}

        if redaction_level == 1:
            # Apply regex-based redaction to the text
            fully_redacted_text = redact_all(text)
        elif redaction_level == 2:
            # Redact entities using NER Method 1
            fully_redacted_text = redact_entities1(text)
        elif redaction_level == 3:
            # Redact entities using NER Method 2
            fully_redacted_text = redact_entities2(text)
        elif redaction_level == 4:
            # Redact entities using NER Method 3
            fully_redacted_text = redact_entities3(text)
        else:
            raise HTTPException(status_code=400, detail="Invalid redaction level. Supported levels: 0, 1, 2, 3, 4.")

        return {"redacted_text": fully_redacted_text}

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"An error occurred: {str(e)}")

@app.post("/redact-pdf/")
async def redact_pdf(file: UploadFile = File(...), redaction_level: int = Form(...)):
    """
    Endpoint to redact sensitive data from a PDF file and return the redacted PDF.
    """
    try:
        # Validate file type
        content_type = file.content_type
        valid_types = [
           "application/pdf",
           "text/plain",
           "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
           "application/vnd.openxmlformats-officedocument.presentationml.presentation",
           "application/octet-stream",  # Assuming logs
           "text/csv" 
        ]

        if content_type not in valid_types:
            raise HTTPException(
                status_code=400,
                detail="Invalid file type. Please upload a PDF, TXT, DOCX, PPTX, or LOG file."
            )

        file_data = await file.read()

        if content_type == "application/pdf":
            return await handle_pdf(file_data, file.filename, redaction_level)

        elif content_type == "text/plain":
            return handle_txt(file_data, file.filename, redaction_level)

        elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return handle_docx(file_data, file.filename, redaction_level)

        elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return handle_pptx(file_data, file.filename, redaction_level)
        
        elif content_type == "text/csv":
            return handle_csv(file_data, file.filename, redaction_level)

        elif content_type == "application/octet-stream":  # Assuming logs are octet-stream
            return handle_log(file_data, file.filename, redaction_level)

    except HTTPException as http_err:
        raise http_err
    except Exception as e:
        logging.error(f"Unexpected error during file redaction: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"An unexpected error occurred: {str(e)}"
        )


async def handle_pdf(file_data, filename, redaction_level):
    try:
        pdf_file = BytesIO(file_data)
        reader = PdfReader(pdf_file)
        _ = reader.pages[0]
        pdf_file.seek(0)

        POPPLER_PATH = r"C:\\Users\\User\\Downloads\\Release-24.08.0-0\\poppler-24.08.0\\Library\\bin"

        # Convert PDF to images
        try:
            images = convert_from_bytes(pdf_file.read(), poppler_path=POPPLER_PATH)
        except Exception as poppler_error:
            raise HTTPException(
                status_code=500,
                detail=f"Error converting PDF to images. Ensure Poppler is installed and configured correctly. {str(poppler_error)}"
            )

        redacted_images = []
        for image in images:
            if redaction_level == 0:
                # Level 0: No redaction
                redacted_images.append(image)
                continue

            ocr_text = pytesseract.image_to_string(image)
            if redaction_level == 1:
                # Level 1: Apply regex-based redaction
                ocr_text = redact_all(ocr_text)
            elif redaction_level == 2:
                # Level 2: Apply NER-based redaction (Method 1)
                ocr_text = redact_entities1(ocr_text)
            elif redaction_level == 3:
                # Level 3: Apply NER-based redaction (Method 2)
                ocr_text = redact_entities2(ocr_text)
            elif redaction_level == 4:
                # Level 4: Apply NER-based redaction (Method 3)
                ocr_text = redact_entities3(ocr_text)

            draw = ImageDraw.Draw(image)
            tsv_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DATAFRAME)
            for _, row in tsv_data.iterrows():
                text = row.get("text", "")
                if isinstance(text, str) and text.strip():
                    if redaction_level == 1:
                        redacted_word = redact_all(text)
                    elif redaction_level == 2:
                        redacted_word = redact_entities1(text)
                    elif redaction_level == 3:
                        redacted_word = redact_entities2(text)
                    elif redaction_level == 4:
                        redacted_word = redact_entities3(text)
                    else:
                        redacted_word = text

                    if redacted_word != text:
                        left, top, width, height = int(row['left']), int(row['top']), int(row['width']), int(row['height'])
                        draw.rectangle([left, top, left + width, top + height], fill="black")

            redacted_images.append(image)

        output_pdf = BytesIO()
        redacted_images[0].save(
            output_pdf,
            format="PDF",
            save_all=True,
            append_images=redacted_images[1:]
        )
        output_pdf.seek(0)

        return StreamingResponse(
            output_pdf,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing PDF: {str(e)}")

def handle_txt(file_data, filename, redaction_level):
    try:
        text = file_data.decode("utf-8", errors="replace")
        if redaction_level == 1:
            text = redact_all(text)
        elif redaction_level == 2:
            text = redact_entities1(text)
        elif redaction_level == 3:
            text = redact_entities2(text)
        elif redaction_level == 4:
            text = redact_entities3(text)

        output_file = BytesIO(text.encode("utf-8"))

        return StreamingResponse(
            output_file,
            media_type="text/plain",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing TXT file: {str(e)}")

def handle_docx(file_data, filename, redaction_level):
    try:
        doc = Document(BytesIO(file_data))
        for paragraph in doc.paragraphs:
            if redaction_level == 1:
                paragraph.text = redact_all(paragraph.text)
            elif redaction_level == 2:
                paragraph.text = redact_entities1(paragraph.text)
            elif redaction_level == 3:
                paragraph.text = redact_entities2(paragraph.text)
            elif redaction_level == 4:
                paragraph.text = redact_entities3(paragraph.text)

        output_docx = BytesIO()
        doc.save(output_docx)
        output_docx.seek(0)

        return StreamingResponse(
            output_docx,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing DOCX file: {str(e)}")

def handle_pptx(file_data, filename, redaction_level):
    try:
        presentation = Presentation(BytesIO(file_data))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if redaction_level == 1:
                                run.text = redact_all(run.text)
                            elif redaction_level == 2:
                                run.text = redact_entities1(run.text)
                            elif redaction_level == 3:
                                run.text = redact_entities2(run.text)
                            elif redaction_level == 4:
                                run.text = redact_entities3(run.text)

        output_pptx = BytesIO()
        presentation.save(output_pptx)
        output_pptx.seek(0)

        return StreamingResponse(
            output_pptx,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing PPTX file: {str(e)}")

def handle_csv(file_data, filename, redaction_level):
    try:
        # Attempt to read the CSV file
        try:
            # Handle encoding errors by catching UnicodeDecodeError explicitly
            data = pd.read_csv(BytesIO(file_data), encoding="utf-8")
        except UnicodeDecodeError:
            raise HTTPException(
                status_code=400,
                detail="Error reading CSV file. Ensure the file is in a valid CSV format and encoded in UTF-8."
            )

        # Function to redact sensitive data in each cell
        def redact_cell(cell):
            if isinstance(cell, str):
                if redaction_level == 1:
                    cell = redact_all(cell)
                elif redaction_level == 2:
                    cell = redact_entities1(cell)
                elif redaction_level == 3:
                    cell = redact_entities2(cell)
                elif redaction_level == 4:
                    cell = redact_entities3(cell)
            return cell

        # Apply redaction to the DataFrame
        try:
            redacted_data = data.applymap(redact_cell)
        except Exception as process_error:
            raise HTTPException(
                status_code=500,
                detail=f"Error processing CSV content. {str(process_error)}"
            )

        # Write the redacted DataFrame back to a CSV buffer
        output_csv = BytesIO()
        try:
            redacted_data.to_csv(output_csv, index=False, encoding="utf-8")
        except Exception as write_error:
            raise HTTPException(
                status_code=500,
                detail=f"Error writing redacted CSV file. {str(write_error)}"
            )

        output_csv.seek(0)

        return StreamingResponse(
            output_csv,
            media_type="text/csv",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Unexpected error processing CSV file: {str(e)}")

def handle_log(file_data, filename, redaction_level):
    try:
        text = file_data.decode("utf-8")
        if redaction_level == 1:
            text = redact_all(text)
        elif redaction_level == 2:
            text = redact_entities1(text)
        elif redaction_level == 3:
            text = redact_entities2(text)
        elif redaction_level == 4:
            text = redact_entities3(text)

        output_file = BytesIO(text.encode("utf-8"))

        return StreamingResponse(
            output_file,
            media_type="text/plain",
            headers={"Content-Disposition": f"attachment; filename=redacted_{filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing LOG file: {str(e)}")






# Exception handler for validation errors
@app.exception_handler(RequestValidationError)
async def validation_exception_handler(request: Request, exc: RequestValidationError):
    logging.error(f"Validation error: {exc}")
    return JSONResponse(
        status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
        content={"detail": exc.errors()},
    )
